"""Editable, dependency-free-of-external-assets campaign presentation export."""

from __future__ import annotations

from collections.abc import Mapping, Sequence
from dataclasses import dataclass, field
from datetime import date, datetime
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)

INK = RGBColor(24, 38, 55)
MUTED = RGBColor(93, 108, 122)
WHITE = RGBColor(255, 255, 255)
DAYLIGHT = RGBColor(247, 250, 248)
MIST = RGBColor(231, 241, 239)
MINT = RGBColor(198, 238, 222)
AURORA = RGBColor(73, 190, 155)
SKY = RGBColor(109, 185, 224)
LILAC = RGBColor(190, 174, 232)
LINE = RGBColor(215, 226, 224)
FONT_FAMILY = "Aptos"


@dataclass(slots=True)
class ContentItem:
    """One editable row in the content matrix."""

    title: str
    channel: str = ""
    format: str = ""
    message: str = ""
    call_to_action: str = ""
    status: str = ""


@dataclass(slots=True)
class ReviewMetadata:
    """Generation provenance and review state shown on the final slide."""

    source_labels: Sequence[str] = field(default_factory=tuple)
    generated_by: str = ""
    generated_at: str | date | datetime = ""
    reviewer: str = ""
    review_status: str = ""
    review_notes: str = ""


@dataclass(slots=True)
class PresentationInput:
    """Plain structured input accepted by :func:`generate_presentation_pptx`."""

    locale: str
    campaign_title: str
    brand: str
    product: str
    audience: str
    objective: str
    core_message: str
    proof_points: Sequence[str] = field(default_factory=tuple)
    content_items: Sequence[ContentItem | Mapping[str, Any]] = field(default_factory=tuple)
    provenance: Sequence[str] = field(default_factory=tuple)
    is_draft: bool = True
    review_metadata: ReviewMetadata | Mapping[str, Any] = field(default_factory=ReviewMetadata)


def generate_presentation_pptx(
    payload: PresentationInput | Mapping[str, Any] | Any,
) -> bytes:
    """Generate a polished, fully editable 16:9 PPTX and return it as bytes.

    ``payload`` may be the dataclass above, a mapping, or a Pydantic-style object
    exposing attributes with the same names. All slide content is made from native
    PowerPoint text, shapes, and tables; no external assets or services are used.
    """

    data = _coerce_input(payload)
    labels = _labels(data.locale)
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    _add_cover(presentation, data, labels)
    _add_summary(presentation, data, labels)
    _add_selling_points(presentation, data, labels)
    _add_content_matrix(presentation, data, labels)
    _add_provenance_review(presentation, data, labels)

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _coerce_input(payload: PresentationInput | Mapping[str, Any] | Any) -> PresentationInput:
    if isinstance(payload, PresentationInput):
        return payload

    review_value = _read(payload, "review_metadata", {})
    review = (
        review_value
        if isinstance(review_value, ReviewMetadata)
        else ReviewMetadata(
            source_labels=_string_sequence(_read(review_value, "source_labels", ())),
            generated_by=_text_value(_read(review_value, "generated_by", "")),
            generated_at=_read(review_value, "generated_at", ""),
            reviewer=_text_value(_read(review_value, "reviewer", "")),
            review_status=_text_value(_read(review_value, "review_status", "")),
            review_notes=_text_value(_read(review_value, "review_notes", "")),
        )
    )
    items = tuple(
        _coerce_content_item(item) for item in (_read(payload, "content_items", ()) or ())
    )
    return PresentationInput(
        locale=_text_value(_read(payload, "locale", "en")),
        campaign_title=_text_value(_read(payload, "campaign_title", "")),
        brand=_text_value(_read(payload, "brand", "")),
        product=_text_value(_read(payload, "product", "")),
        audience=_text_value(_read(payload, "audience", "")),
        objective=_text_value(_read(payload, "objective", "")),
        core_message=_text_value(_read(payload, "core_message", "")),
        proof_points=_string_sequence(_read(payload, "proof_points", ())),
        content_items=items,
        provenance=_string_sequence(_read(payload, "provenance", ())),
        is_draft=bool(_read(payload, "is_draft", True)),
        review_metadata=review,
    )


def _coerce_content_item(value: ContentItem | Mapping[str, Any] | Any) -> ContentItem:
    if isinstance(value, ContentItem):
        return value
    return ContentItem(
        title=_text_value(_read_first(value, ("title", "name"), "")),
        channel=_text_value(_read_first(value, ("channel", "platform"), "")),
        format=_text_value(_read_first(value, ("format", "content_type"), "")),
        message=_text_value(_read_first(value, ("message", "body", "content"), "")),
        call_to_action=_text_value(_read_first(value, ("call_to_action", "cta"), "")),
        status=_text_value(_read(value, "status", "")),
    )


def _read(value: Any, name: str, default: Any) -> Any:
    if isinstance(value, Mapping):
        return value.get(name, default)
    return getattr(value, name, default)


def _read_first(value: Any, names: Sequence[str], default: Any) -> Any:
    for name in names:
        result = _read(value, name, None)
        if result is not None:
            return result
    return default


def _text_value(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (date, datetime)):
        return value.isoformat()
    return str(value)


def _string_sequence(value: Any) -> tuple[str, ...]:
    if value is None:
        return ()
    if isinstance(value, str):
        return (value,)
    return tuple(_text_value(item) for item in value)


def _labels(locale: str) -> dict[str, str]:
    normalized_locale = locale.lower()
    if normalized_locale.startswith(("zh-hk", "zh-tw", "zh-hant")):
        return {
            "deck": "營銷活動提案",
            "summary": "活動摘要",
            "brand": "品牌",
            "product": "產品",
            "audience": "目標受眾",
            "objective": "活動目標",
            "message": "核心信息",
            "selling": "核心賣點",
            "selling_subtitle": "以清晰、可驗證的信息建立購買理由",
            "matrix": "內容矩陣",
            "matrix_subtitle": "跨渠道執行概覽",
            "content": "內容",
            "channel": "渠道",
            "format": "形式",
            "cta_status": "行動 / 狀態",
            "more_items": "另有 {count} 項未顯示",
            "provenance": "來源與審核",
            "sources": "信息來源",
            "generated_by": "生成方",
            "generated_at": "生成時間",
            "reviewer": "審核人",
            "review_status": "審核狀態",
            "notes": "審核備註",
            "not_provided": "未提供",
            "draft_notice": "草稿 · 未經完整審核，請勿直接發佈",
        }
    if normalized_locale.startswith("zh"):
        return {
            "deck": "营销活动提案",
            "summary": "活动摘要",
            "brand": "品牌",
            "product": "产品",
            "audience": "目标受众",
            "objective": "活动目标",
            "message": "核心信息",
            "selling": "核心卖点",
            "selling_subtitle": "以清晰、可验证的信息建立购买理由",
            "matrix": "内容矩阵",
            "matrix_subtitle": "跨渠道执行概览",
            "content": "内容",
            "channel": "渠道",
            "format": "形式",
            "cta_status": "行动 / 状态",
            "more_items": "另有 {count} 项未显示",
            "provenance": "来源与审核",
            "sources": "信息来源",
            "generated_by": "生成方",
            "generated_at": "生成时间",
            "reviewer": "审核人",
            "review_status": "审核状态",
            "notes": "审核备注",
            "not_provided": "未提供",
            "draft_notice": "草稿 · 未经完整审核，请勿直接发布",
        }
    return {
        "deck": "Campaign presentation",
        "summary": "Campaign summary",
        "brand": "Brand",
        "product": "Product",
        "audience": "Audience",
        "objective": "Objective",
        "message": "Core message",
        "selling": "Selling points",
        "selling_subtitle": "Clear, supportable reasons to believe",
        "matrix": "Content matrix",
        "matrix_subtitle": "Cross-channel execution overview",
        "content": "Content",
        "channel": "Channel",
        "format": "Format",
        "cta_status": "CTA / status",
        "more_items": "{count} more items not shown",
        "provenance": "Provenance & review",
        "sources": "Sources",
        "generated_by": "Generated by",
        "generated_at": "Generated at",
        "reviewer": "Reviewer",
        "review_status": "Review status",
        "notes": "Review notes",
        "not_provided": "Not provided",
        "draft_notice": "DRAFT · NOT FULLY REVIEWED · DO NOT PUBLISH",
    }


def _new_slide(presentation: Presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = DAYLIGHT

    _add_circle(slide, 11.6, -0.65, 2.25, MINT)
    _add_circle(slide, 12.15, 0.25, 1.25, SKY)
    _add_circle(slide, -0.55, 6.65, 1.35, LILAC)
    return slide


def _add_circle(slide, x: float, y: float, size: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(size),
        Inches(size),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_cover(presentation: Presentation, data: PresentationInput, labels: dict[str, str]):
    slide = _new_slide(presentation)
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(0.72),
        Inches(1.4),
        Inches(0.13),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = AURORA
    accent.line.fill.background()

    _add_text(
        slide,
        labels["deck"].upper(),
        0.75,
        1.08,
        5.8,
        0.35,
        12,
        MUTED,
        bold=True,
    )
    _add_text(
        slide,
        _truncate(data.campaign_title, 120, labels["deck"]),
        0.72,
        1.62,
        10.7,
        1.55,
        34,
        INK,
        bold=True,
    )
    _add_text(
        slide,
        _truncate(data.core_message, 220, labels["not_provided"]),
        0.76,
        3.38,
        9.7,
        1.15,
        18,
        MUTED,
    )
    if data.is_draft:
        _add_pill(
            slide,
            labels["draft_notice"],
            0.76,
            4.72,
            RGBColor(255, 220, 214),
            width=5.15,
        )

    _add_pill(slide, _truncate(data.brand, 55, labels["brand"]), 0.76, 5.55, MINT)
    _add_pill(slide, _truncate(data.product, 55, labels["product"]), 3.65, 5.55, MIST)
    _add_text(slide, "HEYU AI", 10.55, 6.75, 1.95, 0.28, 11, INK, bold=True, align=PP_ALIGN.RIGHT)
    _add_footer(slide, presentation, 1)
    return slide


def _add_summary(presentation: Presentation, data: PresentationInput, labels: dict[str, str]):
    slide = _new_slide(presentation)
    _add_slide_title(slide, labels["summary"], f"{data.brand} · {data.product}")

    cards = (
        (labels["audience"], data.audience, 0.72, 1.78, 3.72, 1.65, MINT),
        (labels["objective"], data.objective, 4.62, 1.78, 3.72, 1.65, MIST),
        (labels["brand"], data.brand, 8.52, 1.78, 3.72, 1.65, RGBColor(235, 231, 248)),
    )
    for heading, body, x, y, width, height, color in cards:
        _add_card(
            slide,
            heading,
            _truncate(body, 170, labels["not_provided"]),
            x,
            y,
            width,
            height,
            color,
        )

    _add_text(slide, labels["message"].upper(), 0.78, 4.08, 3.0, 0.3, 11, MUTED, bold=True)
    _add_text(
        slide,
        _truncate(data.core_message, 300, labels["not_provided"]),
        0.76,
        4.52,
        11.25,
        1.4,
        24,
        INK,
        bold=True,
    )
    _add_footer(slide, presentation, 2)
    return slide


def _add_selling_points(
    presentation: Presentation,
    data: PresentationInput,
    labels: dict[str, str],
):
    slide = _new_slide(presentation)
    _add_slide_title(slide, labels["selling"], labels["selling_subtitle"])
    points = list(data.proof_points[:6])
    if not points:
        points = [labels["not_provided"]]

    for index, point in enumerate(points):
        column = index % 2
        row = index // 2
        x = 0.76 + column * 6.05
        y = 1.72 + row * 1.62
        number = f"{index + 1:02d}"
        _add_text(slide, number, x, y + 0.03, 0.56, 0.42, 13, AURORA, bold=True)
        _add_text(
            slide,
            _truncate(point, 180),
            x + 0.64,
            y,
            4.92,
            1.03,
            17,
            INK,
            bold=True,
        )
        divider = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x + 0.64),
            Inches(y + 1.18),
            Inches(4.92),
            Inches(0.015),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = LINE
        divider.line.fill.background()

    if len(data.proof_points) > 6:
        _add_text(
            slide,
            labels["more_items"].format(count=len(data.proof_points) - 6),
            0.78,
            6.53,
            5.5,
            0.25,
            10,
            MUTED,
        )
    _add_footer(slide, presentation, 3)
    return slide


def _add_content_matrix(
    presentation: Presentation,
    data: PresentationInput,
    labels: dict[str, str],
):
    slide = _new_slide(presentation)
    _add_slide_title(slide, labels["matrix"], labels["matrix_subtitle"])
    shown_items = list(data.content_items[:7])
    row_count = max(1, len(shown_items)) + 1
    table_shape = slide.shapes.add_table(
        row_count,
        4,
        Inches(0.72),
        Inches(1.67),
        Inches(11.88),
        Inches(4.85),
    )
    table = table_shape.table
    widths = (4.45, 1.95, 1.95, 3.53)
    for column, width in zip(table.columns, widths, strict=True):
        column.width = Inches(width)

    headers = (labels["content"], labels["channel"], labels["format"], labels["cta_status"])
    for column_index, header in enumerate(headers):
        _style_cell(table.cell(0, column_index), header, AURORA, WHITE, 11, bold=True)

    if not shown_items:
        shown_items = [ContentItem(title=labels["not_provided"])]

    for row_index, item in enumerate(shown_items, start=1):
        title_and_message = _truncate(item.title, 80, labels["not_provided"])
        if item.message:
            title_and_message += "\n" + _truncate(item.message, 130)
        cta_status = " · ".join(
            part
            for part in (_truncate(item.call_to_action, 70), _truncate(item.status, 35))
            if part
        )
        values = (
            title_and_message,
            _truncate(item.channel, 55, "—"),
            _truncate(item.format, 55, "—"),
            cta_status or "—",
        )
        fill = WHITE if row_index % 2 else RGBColor(241, 247, 245)
        for column_index, value in enumerate(values):
            _style_cell(table.cell(row_index, column_index), value, fill, INK, 10)

    if len(data.content_items) > 7:
        _add_text(
            slide,
            labels["more_items"].format(count=len(data.content_items) - 7),
            8.15,
            6.57,
            4.0,
            0.24,
            10,
            MUTED,
            align=PP_ALIGN.RIGHT,
        )
    _add_footer(slide, presentation, 4)
    return slide


def _add_provenance_review(
    presentation: Presentation,
    data: PresentationInput,
    labels: dict[str, str],
):
    slide = _new_slide(presentation)
    _add_slide_title(slide, labels["provenance"], data.campaign_title)
    review = data.review_metadata
    source_values = tuple(
        dict.fromkeys(source for source in (*data.provenance, *review.source_labels) if source)
    )
    source_text = _truncate(
        "\n".join(f"• {_truncate(source, 120)}" for source in source_values[:8]),
        680,
        labels["not_provided"],
    )

    _add_card(
        slide,
        labels["sources"],
        source_text,
        0.72,
        1.68,
        5.82,
        4.86,
        MIST,
        body_size=12,
    )

    metadata = (
        (labels["generated_by"], review.generated_by),
        (labels["generated_at"], _text_value(review.generated_at)),
        (labels["reviewer"], review.reviewer),
        (labels["review_status"], review.review_status),
    )
    for index, (heading, value) in enumerate(metadata):
        x = 6.78 + (index % 2) * 2.82
        y = 1.68 + (index // 2) * 1.32
        _add_card(
            slide,
            heading,
            _truncate(value, 80, labels["not_provided"]),
            x,
            y,
            2.57,
            1.08,
            WHITE,
            body_size=12,
        )

    _add_card(
        slide,
        labels["notes"],
        _truncate(review.review_notes, 380, labels["not_provided"]),
        6.78,
        4.35,
        5.64,
        2.19,
        RGBColor(235, 231, 248),
        body_size=12,
    )
    _add_footer(slide, presentation, 5)
    return slide


def _add_slide_title(slide, title: str, subtitle: str) -> None:
    _add_text(slide, _truncate(title, 90), 0.72, 0.54, 8.7, 0.55, 27, INK, bold=True)
    _add_text(slide, _truncate(subtitle, 140), 0.74, 1.13, 9.8, 0.28, 11, MUTED)


def _add_card(
    slide,
    heading: str,
    body: str,
    x: float,
    y: float,
    width: float,
    height: float,
    fill_color: RGBColor,
    *,
    body_size: int = 14,
) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = LINE
    card.line.width = Pt(0.8)
    _add_text(slide, heading.upper(), x + 0.22, y + 0.18, width - 0.44, 0.26, 10, MUTED, bold=True)
    _add_text(
        slide,
        body,
        x + 0.22,
        y + 0.55,
        width - 0.44,
        height - 0.69,
        body_size,
        INK,
        bold=height <= 1.7,
    )


def _add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    color: RGBColor,
    *,
    width: float = 2.62,
) -> None:
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(0.57),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    frame = pill.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.16)
    frame.margin_right = Inches(0.16)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = INK


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    size: int,
    color: RGBColor,
    *,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = 0
    run = paragraph.add_run()
    run.text = _clean_text(text)
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _style_cell(
    cell,
    text: str,
    fill_color: RGBColor,
    font_color: RGBColor,
    size: int,
    *,
    bold: bool = False,
) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_color
    cell.margin_left = Inches(0.12)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.07)
    cell.margin_bottom = Inches(0.05)
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.space_after = 0
    run = paragraph.add_run()
    run.text = _clean_text(text)
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = font_color


def _add_footer(slide, presentation: Presentation, slide_number: int) -> None:
    _add_text(slide, f"{slide_number:02d}", 12.05, 7.02, 0.5, 0.2, 9, MUTED, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.72),
        Inches(7.08),
        Inches(10.95),
        Inches(0.012),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def _clean_text(value: str) -> str:
    raw = _text_value(value).replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    return "\n".join(" ".join(line.split()) for line in raw.split("\n"))


def _truncate(value: Any, limit: int, fallback: str = "") -> str:
    text = _clean_text(_text_value(value))
    if not text:
        return fallback
    if len(text) <= limit:
        return text
    return text[: max(1, limit - 1)].rstrip() + "…"


__all__ = [
    "ContentItem",
    "PresentationInput",
    "ReviewMetadata",
    "generate_presentation_pptx",
]
