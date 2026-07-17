from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from itertools import islice
from pathlib import PurePath
from typing import Any, Literal
from zipfile import BadZipFile, ZipFile

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from pypdf.errors import PdfReadError

PDF_MEDIA_TYPE = "application/pdf"
PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
DOCX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
DEFAULT_MAX_PAGES = 100
DEFAULT_MAX_SLIDES = 100
DEFAULT_MAX_PARAGRAPHS = 1_000
DEFAULT_MAX_CHARACTERS = 100_000
MAX_OFFICE_MEMBERS = 2_000
MAX_OFFICE_UNCOMPRESSED_BYTES = 100 * 1024 * 1024
MAX_OFFICE_COMPRESSION_RATIO = 200

FragmentKind = Literal["page", "slide", "paragraph"]
DocumentKind = Literal["pdf", "pptx", "docx"]


@dataclass(frozen=True)
class DocumentFragment:
    kind: FragmentKind
    number: int
    text: str


@dataclass(frozen=True)
class DocumentExtractionResult:
    document_kind: DocumentKind
    full_text: str
    fragments: tuple[DocumentFragment, ...]
    warnings: tuple[str, ...]


class DocumentImportError(ValueError):
    """Base error with fields suitable for a FastAPI HTTPException."""

    status_code = 422
    code = "document_import_error"

    def __init__(self, detail: str) -> None:
        super().__init__(detail)
        self.detail = detail


class EmptyDocumentError(DocumentImportError):
    status_code = 400
    code = "empty_document"


class UnsupportedDocumentTypeError(DocumentImportError):
    status_code = 415
    code = "unsupported_document_type"


class EncryptedDocumentError(DocumentImportError):
    status_code = 422
    code = "encrypted_document"


class InvalidDocumentError(DocumentImportError):
    status_code = 422
    code = "invalid_document"


class InvalidDocumentLimitError(DocumentImportError):
    status_code = 400
    code = "invalid_document_limit"


def extract_document_text(
    data: bytes,
    *,
    media_type: str | None = None,
    filename: str | None = None,
    max_pages: int = DEFAULT_MAX_PAGES,
    max_slides: int = DEFAULT_MAX_SLIDES,
    max_paragraphs: int = DEFAULT_MAX_PARAGRAPHS,
    max_characters: int = DEFAULT_MAX_CHARACTERS,
) -> DocumentExtractionResult:
    """Extract text from PDF, PPTX, or DOCX bytes without filesystem or network access."""
    if not data:
        raise EmptyDocumentError("The uploaded document is empty.")
    _validate_limit("max_pages", max_pages)
    _validate_limit("max_slides", max_slides)
    _validate_limit("max_paragraphs", max_paragraphs)
    _validate_limit("max_characters", max_characters)

    document_kind = _detect_document_kind(data, media_type=media_type, filename=filename)
    if document_kind == "pdf":
        raw_fragments, warnings = _extract_pdf(data, max_pages=max_pages)
        fragment_kind: FragmentKind = "page"
    elif document_kind == "pptx":
        raw_fragments, warnings = _extract_pptx(data, max_slides=max_slides)
        fragment_kind = "slide"
    else:
        raw_fragments, warnings = _extract_docx(data, max_paragraphs=max_paragraphs)
        fragment_kind = "paragraph"

    fragments, full_text, was_truncated = _truncate_fragments(
        raw_fragments,
        kind=fragment_kind,
        max_characters=max_characters,
    )
    if was_truncated:
        warnings.append(f"Extracted text was truncated at the {max_characters}-character limit.")
    if not full_text:
        warnings.append("The document contains no extractable text.")
        if document_kind == "pdf":
            warnings.append(
                "This may be a scanned PDF. OCR is optional and is not configured in "
                "the zero-cost local demo; export a searchable PDF or paste reviewed text."
            )

    return DocumentExtractionResult(
        document_kind=document_kind,
        full_text=full_text,
        fragments=tuple(fragments),
        warnings=tuple(warnings),
    )


def _validate_limit(name: str, value: int) -> None:
    if isinstance(value, bool) or not isinstance(value, int) or value < 1:
        raise InvalidDocumentLimitError(f"{name} must be a positive integer.")


def _detect_document_kind(
    data: bytes,
    *,
    media_type: str | None,
    filename: str | None,
) -> DocumentKind:
    normalized_media_type = (media_type or "").partition(";")[0].strip().lower()
    if normalized_media_type == PDF_MEDIA_TYPE:
        return "pdf"
    if normalized_media_type == PPTX_MEDIA_TYPE:
        return "pptx"
    if normalized_media_type == DOCX_MEDIA_TYPE:
        return "docx"

    suffix = PurePath(filename or "").suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix == ".pptx":
        return "pptx"
    if suffix == ".docx":
        return "docx"

    if data.startswith(b"%PDF-"):
        return "pdf"
    if data.startswith(b"PK\x03\x04"):
        return _detect_office_archive_kind(data)

    raise UnsupportedDocumentTypeError("Only PDF, PPTX, and DOCX documents are supported.")


def _extract_pdf(data: bytes, *, max_pages: int) -> tuple[list[str], list[str]]:
    try:
        reader = PdfReader(BytesIO(data), strict=False)
    except (PdfReadError, OSError, ValueError, TypeError) as exc:
        raise InvalidDocumentError("The PDF document is corrupt or invalid.") from exc

    if reader.is_encrypted:
        raise EncryptedDocumentError("Encrypted PDF documents are not supported.")

    try:
        page_count = len(reader.pages)
    except (PdfReadError, OSError, ValueError, TypeError) as exc:
        raise InvalidDocumentError("The PDF document is corrupt or invalid.") from exc

    warnings: list[str] = []
    if page_count > max_pages:
        warnings.append(f"Only the first {max_pages} of {page_count} PDF pages were processed.")

    page_texts: list[str] = []
    for page_number, page in enumerate(reader.pages[:max_pages], start=1):
        try:
            page_texts.append((page.extract_text() or "").strip())
        except Exception as exc:
            page_texts.append("")
            warnings.append(f"Text could not be extracted from PDF page {page_number}: {exc}.")
    return page_texts, warnings


def _extract_pptx(data: bytes, *, max_slides: int) -> tuple[list[str], list[str]]:
    _validate_office_archive(data, kind="pptx")
    try:
        presentation = Presentation(BytesIO(data))
        slide_count = len(presentation.slides)
        slide_texts = [
            _extract_slide_text(slide) for slide in islice(presentation.slides, max_slides)
        ]
    except (OSError, ValueError, KeyError, TypeError) as exc:
        raise InvalidDocumentError("The PPTX document is corrupt or invalid.") from exc
    except Exception as exc:
        raise InvalidDocumentError("The PPTX document could not be read.") from exc

    warnings: list[str] = []
    if slide_count > max_slides:
        warnings.append(f"Only the first {max_slides} of {slide_count} PPTX slides were processed.")
    return slide_texts, warnings


def _extract_docx(data: bytes, *, max_paragraphs: int) -> tuple[list[str], list[str]]:
    _validate_office_archive(data, kind="docx")
    try:
        document = Document(BytesIO(data))
        all_fragments = [
            text
            for text in (
                [paragraph.text.strip() for paragraph in document.paragraphs]
                + [
                    " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    for table in document.tables
                    for row in table.rows
                ]
            )
            if text
        ]
    except (OSError, ValueError, KeyError, TypeError) as exc:
        raise InvalidDocumentError("The DOCX document is corrupt or invalid.") from exc
    except Exception as exc:
        raise InvalidDocumentError("The DOCX document could not be read.") from exc

    warnings: list[str] = []
    if len(all_fragments) > max_paragraphs:
        warnings.append(
            f"Only the first {max_paragraphs} of {len(all_fragments)} DOCX text "
            "blocks were processed."
        )
    return all_fragments[:max_paragraphs], warnings


def _detect_office_archive_kind(data: bytes) -> Literal["pptx", "docx"]:
    try:
        with ZipFile(BytesIO(data)) as archive:
            names = {member.filename for member in archive.infolist()}
    except (BadZipFile, OSError, ValueError) as exc:
        raise InvalidDocumentError("The Office document is corrupt or invalid.") from exc
    if "ppt/presentation.xml" in names:
        return "pptx"
    if "word/document.xml" in names:
        return "docx"
    raise UnsupportedDocumentTypeError("The uploaded archive is not a valid PPTX or DOCX file.")


def _validate_office_archive(data: bytes, *, kind: Literal["pptx", "docx"]) -> None:
    try:
        with ZipFile(BytesIO(data)) as archive:
            members = archive.infolist()
    except (BadZipFile, OSError, ValueError) as exc:
        raise InvalidDocumentError("The PPTX document is corrupt or invalid.") from exc

    label = kind.upper()
    if len(members) > MAX_OFFICE_MEMBERS:
        raise InvalidDocumentError(f"The {label} document contains too many archive entries.")

    names = {member.filename for member in members}
    required_part = "ppt/presentation.xml" if kind == "pptx" else "word/document.xml"
    if "[Content_Types].xml" not in names or required_part not in names:
        raise InvalidDocumentError(f"The uploaded archive is not a valid {label} document.")

    total_uncompressed = sum(member.file_size for member in members)
    if total_uncompressed > MAX_OFFICE_UNCOMPRESSED_BYTES:
        raise InvalidDocumentError(
            f"The {label} document expands beyond the safe processing limit."
        )

    for member in members:
        if (
            member.file_size > 1_000_000
            and member.compress_size > 0
            and member.file_size / member.compress_size > MAX_OFFICE_COMPRESSION_RATIO
        ):
            raise InvalidDocumentError(f"The {label} document has an unsafe compression ratio.")


def _extract_slide_text(slide: Any) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        parts.extend(_extract_shape_text(shape))
    return "\n".join(part for part in parts if part).strip()


def _extract_shape_text(shape: Any) -> list[str]:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        parts: list[str] = []
        for child in shape.shapes:
            parts.extend(_extract_shape_text(child))
        return parts

    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        return [text] if text else []

    if getattr(shape, "has_table", False):
        return [
            cell.text.strip() for row in shape.table.rows for cell in row.cells if cell.text.strip()
        ]

    return []


def _truncate_fragments(
    texts: list[str],
    *,
    kind: FragmentKind,
    max_characters: int,
) -> tuple[list[DocumentFragment], str, bool]:
    fragments: list[DocumentFragment] = []
    used_characters = 0
    has_text = False
    truncated = False

    for number, text in enumerate(texts, start=1):
        separator_length = 2 if has_text and text else 0
        available_for_text = max_characters - used_characters - separator_length
        if text and available_for_text <= 0:
            truncated = True
            break

        fragment_text = text[:available_for_text]
        fragments.append(DocumentFragment(kind=kind, number=number, text=fragment_text))
        used_characters += separator_length + len(fragment_text)
        has_text = has_text or bool(fragment_text)

        if len(fragment_text) < len(text):
            truncated = True
            break

    if len(fragments) < len(texts) and any(texts[len(fragments) :]):
        truncated = True

    full_text = "\n\n".join(fragment.text for fragment in fragments if fragment.text)
    return fragments, full_text, truncated
