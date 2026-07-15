from io import BytesIO

import pytest
from pptx import Presentation

from app.presentation_export import (
    ContentItem,
    PresentationInput,
    ReviewMetadata,
    generate_presentation_pptx,
)


def _presentation(payload) -> Presentation:
    output = generate_presentation_pptx(payload)
    assert output.startswith(b"PK")
    return Presentation(BytesIO(output))


def _slide_text(slide) -> str:
    values: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            values.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                values.extend(cell.text for cell in row.cells)
    return "\n".join(values)


def test_generates_editable_wide_campaign_deck_with_expected_content():
    payload = PresentationInput(
        locale="en-US",
        campaign_title="Sunrise Harvest Launch",
        brand="Heyu Fields",
        product="Aurora Tomatoes",
        audience="Urban families who value transparent sourcing",
        objective="Build trust and drive first purchase",
        core_message="Traceable freshness from field to family table.",
        proof_points=(
            "Picked at peak ripeness",
            "Origin and handling records are available",
            "Cold-chain delivery protects freshness",
        ),
        content_items=(
            ContentItem(
                title="Meet the grower",
                channel="Short video",
                format="30-second story",
                message="Show the harvest and the people behind it.",
                call_to_action="Explore the harvest",
                status="Approved",
            ),
            ContentItem(
                title="Freshness explainer",
                channel="Social",
                format="Carousel",
                message="Explain picking and delivery timing.",
                call_to_action="Learn more",
                status="In review",
            ),
        ),
        provenance=("Farm visit notes", "Approved product specification"),
        is_draft=False,
        review_metadata=ReviewMetadata(
            source_labels=("Fulfilment confirmation",),
            generated_by="Heyu AI",
            generated_at="2026-07-14",
            reviewer="Campaign review team",
            review_status="Approved",
            review_notes="Product and delivery claims checked against current records.",
        ),
    )

    presentation = _presentation(payload)

    assert len(presentation.slides) == 5
    assert presentation.slide_width / presentation.slide_height == pytest.approx(16 / 9)

    deck_text = "\n".join(_slide_text(slide) for slide in presentation.slides)
    assert "Sunrise Harvest Launch" in deck_text
    assert "Traceable freshness from field to family table." in deck_text
    assert "Meet the grower" in deck_text
    assert "Farm visit notes" in deck_text
    assert "Campaign review team" in deck_text
    assert "DO NOT PUBLISH" not in deck_text

    assert any(shape.has_table for shape in presentation.slides[3].shapes)
    assert any(
        shape.has_text_frame and shape.text == "Sunrise Harvest Launch"
        for shape in presentation.slides[0].shapes
    )


def test_long_mapping_content_is_truncated_and_kept_inside_slide_bounds():
    long_value = "Long campaign detail " * 500
    payload = {
        "locale": "zh-CN",
        "campaign_title": "超长活动标题" + long_value,
        "brand": "和予品牌" + long_value,
        "product": "极光番茄" + long_value,
        "audience": long_value,
        "objective": long_value,
        "core_message": "可追溯的新鲜承诺 " + long_value,
        "proof_points": [f"卖点 {index} {long_value}" for index in range(20)],
        "content_items": [
            {
                "title": f"内容 {index} {long_value}",
                "channel": long_value,
                "format": long_value,
                "message": long_value,
                "call_to_action": long_value,
                "status": long_value,
            }
            for index in range(20)
        ],
        "provenance": [long_value for _ in range(20)],
        "review_metadata": {
            "source_labels": [long_value for _ in range(20)],
            "generated_by": long_value,
            "generated_at": "2026-07-14T10:30:00+08:00",
            "reviewer": long_value,
            "review_status": long_value,
            "review_notes": long_value,
        },
    }

    presentation = _presentation(payload)

    assert len(presentation.slides) == 5
    assert "可追溯的新鲜承诺" in _slide_text(presentation.slides[0])
    assert "另有 14 项未显示" in _slide_text(presentation.slides[2])
    assert "另有 13 项未显示" in _slide_text(presentation.slides[3])

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                assert len(shape.text) <= 700
                assert shape.left >= 0
                assert shape.top >= 0
                assert shape.left + shape.width <= presentation.slide_width
                assert shape.top + shape.height <= presentation.slide_height
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        assert len(cell.text) <= 220


def test_hong_kong_chinese_uses_traditional_labels_and_deduplicates_sources():
    presentation = _presentation(
        PresentationInput(
            locale="zh-HK",
            campaign_title="時令番茄推廣",
            brand="禾語田園",
            product="日光番茄",
            audience="重視產地資料的家庭",
            objective="清楚說明產品資料",
            core_message="由田間到餐桌，每項信息均可核對。",
            provenance=("已審核產品資料",),
            review_metadata=ReviewMetadata(
                source_labels=("已審核產品資料",),
                review_status="已審核",
            ),
        )
    )

    deck_text = "\n".join(_slide_text(slide) for slide in presentation.slides)
    assert "營銷活動提案" in deck_text
    assert "內容矩陣" in deck_text
    assert "來源與審核" in deck_text
    assert "营销活动提案" not in deck_text
    assert deck_text.count("已審核產品資料") == 1
    assert "未經完整審核" in deck_text
