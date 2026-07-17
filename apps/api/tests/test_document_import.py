from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from pypdf.generic import (
    DecodedStreamObject,
    DictionaryObject,
    NameObject,
    TextStringObject,
)

from app.document_import import (
    DOCX_MEDIA_TYPE,
    PDF_MEDIA_TYPE,
    PPTX_MEDIA_TYPE,
    DocumentImportError,
    EmptyDocumentError,
    EncryptedDocumentError,
    InvalidDocumentError,
    extract_document_text,
)


def make_pdf(*page_texts: str, password: str | None = None) -> bytes:
    writer = PdfWriter()
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_reference = writer._add_object(font)

    for text in page_texts:
        page = writer.add_blank_page(width=612, height=792)
        page[NameObject("/Resources")] = DictionaryObject(
            {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_reference})}
        )
        content = DecodedStreamObject()
        encoded_text = TextStringObject(text).write_to_stream
        text_buffer = BytesIO()
        encoded_text(text_buffer)
        content.set_data(b"BT /F1 12 Tf 72 720 Td " + text_buffer.getvalue() + b" Tj ET")
        page[NameObject("/Contents")] = writer._add_object(content)

    if password:
        writer.encrypt(password)

    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def make_pptx(*slide_texts: str) -> bytes:
    presentation = Presentation()
    presentation.slides._sldIdLst.clear()
    for text in slide_texts:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        text_box.text = text

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def make_docx(*paragraph_texts: str) -> bytes:
    document = Document()
    for text in paragraph_texts:
        document.add_paragraph(text)
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Harvest date"
    table.cell(0, 1).text = "2026-07-17"
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def test_extracts_structured_pdf_text() -> None:
    result = extract_document_text(
        make_pdf("First PDF page", "Second PDF page"),
        media_type=PDF_MEDIA_TYPE,
    )

    assert result.document_kind == "pdf"
    assert "First PDF page" in result.full_text
    assert [fragment.kind for fragment in result.fragments] == ["page", "page"]
    assert [fragment.number for fragment in result.fragments] == [1, 2]
    assert result.warnings == ()


def test_extracts_structured_pptx_text() -> None:
    result = extract_document_text(
        make_pptx("First slide", "Second slide"),
        media_type=PPTX_MEDIA_TYPE,
    )

    assert result.document_kind == "pptx"
    assert "First slide" in result.full_text
    assert [fragment.kind for fragment in result.fragments] == ["slide", "slide"]
    assert [fragment.number for fragment in result.fragments] == [1, 2]
    assert result.warnings == ()


def test_extracts_structured_docx_text_and_table_rows() -> None:
    result = extract_document_text(
        make_docx("Seasonal tomato", "Picked after natural ripening"),
        media_type=DOCX_MEDIA_TYPE,
    )

    assert result.document_kind == "docx"
    assert "Seasonal tomato" in result.full_text
    assert "Harvest date | 2026-07-17" in result.full_text
    assert all(fragment.kind == "paragraph" for fragment in result.fragments)
    assert [fragment.number for fragment in result.fragments] == [1, 2, 3]
    assert result.warnings == ()


def test_empty_input_is_a_client_error() -> None:
    with pytest.raises(EmptyDocumentError) as raised:
        extract_document_text(b"", media_type=PDF_MEDIA_TYPE)

    assert isinstance(raised.value, DocumentImportError)
    assert raised.value.status_code == 400
    assert raised.value.code == "empty_document"


def test_valid_blank_document_returns_warning() -> None:
    result = extract_document_text(make_pdf(""), media_type=PDF_MEDIA_TYPE)

    assert result.full_text == ""
    assert result.fragments[0].text == ""
    assert any("no extractable text" in warning.lower() for warning in result.warnings)


@pytest.mark.parametrize("media_type", [PDF_MEDIA_TYPE, PPTX_MEDIA_TYPE, DOCX_MEDIA_TYPE])
def test_corrupt_document_is_a_client_error(media_type: str) -> None:
    with pytest.raises(InvalidDocumentError) as raised:
        extract_document_text(b"not a real document", media_type=media_type)

    assert raised.value.status_code == 422
    assert raised.value.code == "invalid_document"


def test_encrypted_pdf_is_rejected() -> None:
    with pytest.raises(EncryptedDocumentError) as raised:
        extract_document_text(make_pdf("secret", password="password"), media_type=PDF_MEDIA_TYPE)

    assert raised.value.status_code == 422
    assert raised.value.code == "encrypted_document"


def test_blank_pdf_explains_optional_ocr_degradation() -> None:
    result = extract_document_text(make_pdf(""), media_type=PDF_MEDIA_TYPE)

    assert any("scanned PDF" in warning for warning in result.warnings)
    assert any("OCR is optional" in warning for warning in result.warnings)


def test_character_and_slide_limits_truncate_with_warnings() -> None:
    data = make_pptx("abcdefghij", "second slide")

    character_limited = extract_document_text(
        data,
        media_type=PPTX_MEDIA_TYPE,
        max_characters=5,
    )
    assert character_limited.full_text == "abcde"
    assert character_limited.fragments[0].text == "abcde"
    assert len(character_limited.full_text) == 5
    assert any("character limit" in warning for warning in character_limited.warnings)

    slide_limited = extract_document_text(
        data,
        media_type=PPTX_MEDIA_TYPE,
        max_slides=1,
    )
    assert len(slide_limited.fragments) == 1
    assert slide_limited.full_text == "abcdefghij"
    assert any("first 1 of 2" in warning for warning in slide_limited.warnings)


def test_rejects_pptx_archive_with_unsafe_compression_ratio() -> None:
    output = BytesIO()
    with ZipFile(output, "w", ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "<Types />")
        archive.writestr("ppt/presentation.xml", "<presentation />")
        archive.writestr("ppt/media/oversized.bin", b"A" * 1_100_000)

    with pytest.raises(InvalidDocumentError, match="unsafe compression ratio"):
        extract_document_text(output.getvalue(), media_type=PPTX_MEDIA_TYPE)
