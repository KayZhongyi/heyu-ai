from io import BytesIO

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from sqlalchemy.orm import Session

from app.main import MAX_DOCUMENT_UPLOAD_BYTES
from app.models import ContentVersion, ReviewStatus
from tests.conftest import bootstrap, invite_and_accept
from tests.test_campaign_packages import (
    approve_campaign_assets,
    campaign_payload,
    create_approved_supply,
    create_assets,
)


def make_pdf() -> bytes:
    output = BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    writer.write(output)
    return output.getvalue()


def make_pptx(text: str = "Traceable harvest story") -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text = text
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_document_preview_requires_write_role(
    client: TestClient,
    auth: dict[str, str],
) -> None:
    _, viewer = invite_and_accept(
        client,
        auth,
        "viewer-document@example.com",
        "viewer",
        "viewer-password",
    )
    viewer_auth = {"Authorization": f"Bearer {viewer['access_token']}"}

    unauthorized = client.post(
        "/v1/document-imports/preview",
        files={"file": ("source.pptx", make_pptx())},
    )
    forbidden = client.post(
        "/v1/document-imports/preview",
        headers=viewer_auth,
        files={"file": ("source.pptx", make_pptx())},
    )

    assert unauthorized.status_code == 401
    assert forbidden.status_code == 403


def test_pptx_document_preview_returns_editable_text_and_provenance(
    client: TestClient,
    auth: dict[str, str],
) -> None:
    response = client.post(
        "/v1/document-imports/preview",
        headers=auth,
        files={
            "file": (
                "field-notes.pptx",
                make_pptx(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert response.status_code == 200, response.text
    payload = response.json()
    assert payload["filename"] == "field-notes.pptx"
    assert payload["text"] == "Traceable harvest story"
    assert payload["sections"] == [
        {
            "kind": "slide",
            "number": 1,
            "label": "Slide 1",
            "text": "Traceable harvest story",
        }
    ]
    assert len(payload["content_sha256"]) == 64


def test_document_preview_reports_detected_type_when_filename_has_no_extension(
    client: TestClient,
    auth: dict[str, str],
) -> None:
    response = client.post(
        "/v1/document-imports/preview",
        headers=auth,
        files={"file": ("uploaded-document", make_pdf(), "application/octet-stream")},
    )

    assert response.status_code == 200, response.text
    assert response.json()["media_type"] == "application/pdf"


def test_document_preview_rejects_unsupported_corrupt_and_large_files(
    client: TestClient,
    auth: dict[str, str],
) -> None:
    unsupported = client.post(
        "/v1/document-imports/preview",
        headers=auth,
        files={"file": ("source.docx", b"not-supported")},
    )
    corrupt = client.post(
        "/v1/document-imports/preview",
        headers=auth,
        files={"file": ("source.pdf", b"%PDF-not-valid", "application/pdf")},
    )
    too_large = client.post(
        "/v1/document-imports/preview",
        headers=auth,
        files={"file": ("source.pdf", b"x" * (MAX_DOCUMENT_UPLOAD_BYTES + 1))},
    )

    assert unsupported.status_code == 415
    assert corrupt.status_code == 422
    assert too_large.status_code == 413


def test_pdf_document_preview_reports_blank_scanned_style_document(
    client: TestClient,
    auth: dict[str, str],
) -> None:
    response = client.post(
        "/v1/document-imports/preview",
        headers=auth,
        files={"file": ("blank.pdf", make_pdf(), "application/pdf")},
    )

    assert response.status_code == 200, response.text
    assert response.json()["text"] == ""
    assert any("no extractable text" in item.lower() for item in response.json()["warnings"])


def test_campaign_presentation_download_is_editable_and_tenant_scoped(
    client: TestClient,
    auth: dict[str, str],
) -> None:
    brand, product = create_assets(client, auth)
    campaign_response = client.post(
        "/v1/campaign-packages",
        headers=auth,
        json=campaign_payload(brand, product),
    )
    assert campaign_response.status_code == 201, campaign_response.text
    campaign = campaign_response.json()

    download = client.get(
        f"/v1/campaign-packages/{campaign['id']}/presentation",
        headers=auth,
    )

    assert download.status_code == 200, download.text
    assert download.content.startswith(b"PK")
    assert (
        download.headers["content-type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    presentation = Presentation(BytesIO(download.content))
    assert len(presentation.slides) == 5
    deck_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert campaign["title"] in deck_text
    assert brand["name"] in deck_text
    assert product["name"] in deck_text
    assert "草稿" in deck_text
    assert "请勿直接发布" in deck_text

    second = bootstrap(client, "presentation-second", "presentation-second@example.com")
    second_auth = {"Authorization": f"Bearer {second['access_token']}"}
    hidden = client.get(
        f"/v1/campaign-packages/{campaign['id']}/presentation",
        headers=second_auth,
    )
    assert hidden.status_code == 404


def test_campaign_presentation_marks_stale_approved_content_as_draft(
    client: TestClient,
    auth: dict[str, str],
    db: Session,
) -> None:
    brand, product = create_assets(client, auth)
    approve_campaign_assets(client, auth, brand, product)
    payload = campaign_payload(brand, product)
    payload["create_default_items"] = True
    campaign_response = client.post(
        "/v1/campaign-packages",
        headers=auth,
        json=payload,
    )
    assert campaign_response.status_code == 201, campaign_response.text
    campaign = campaign_response.json()
    supply = create_approved_supply(client, auth, campaign, brand, product)

    for item in campaign["items"]:
        db.add(
            ContentVersion(
                organization_id=campaign["organization_id"],
                project_id=item["project"]["id"],
                brief_revision_id=campaign["current_brief_revision"]["id"],
                supply_snapshot_id=supply["id"],
                version_number=1,
                content={
                    "headline": "Verified seasonal produce",
                    "body": "Approved campaign content.",
                },
                status=ReviewStatus.approved,
                created_by=campaign["created_by"],
                reviewed_by=campaign["created_by"],
                review_note="Current evidence verified",
            )
        )
    db.commit()

    activated = client.patch(
        f"/v1/campaign-packages/{campaign['id']}/status",
        headers=auth,
        json={"status": "active"},
    )
    assert activated.status_code == 200, activated.text

    current_download = client.get(
        f"/v1/campaign-packages/{campaign['id']}/presentation",
        headers=auth,
    )
    assert current_download.status_code == 200, current_download.text
    current_deck = Presentation(BytesIO(current_download.content))
    current_text = "\n".join(
        shape.text
        for slide in current_deck.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "请勿直接发布" not in current_text

    create_approved_supply(client, auth, campaign, brand, product)
    stale_campaign = client.get(
        f"/v1/campaign-packages/{campaign['id']}",
        headers=auth,
    )
    assert stale_campaign.status_code == 200, stale_campaign.text
    assert all(item["approved_version_id"] is None for item in stale_campaign.json()["items"]), (
        stale_campaign.json()
    )

    stale_download = client.get(
        f"/v1/campaign-packages/{campaign['id']}/presentation",
        headers=auth,
    )
    assert stale_download.status_code == 200, stale_download.text
    stale_deck = Presentation(BytesIO(stale_download.content))
    stale_text = "\n".join(
        shape.text
        for slide in stale_deck.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "请勿直接发布" in stale_text


def test_campaign_presentation_marks_unmapped_approved_claims_as_draft(
    client: TestClient,
    auth: dict[str, str],
    db: Session,
) -> None:
    brand, product = create_assets(client, auth)
    initial_product_payload = {
        key: product[key]
        for key in (
            "brand_id",
            "name",
            "origin",
            "specification",
            "price_display",
            "shelf_life",
            "storage_method",
            "selling_points",
            "prohibited_claims",
        )
    }
    initial_product_payload["origin"] = "Yunnan"
    updated_product = client.put(
        f"/v1/products/{product['id']}",
        headers=auth,
        json=initial_product_payload,
    )
    assert updated_product.status_code == 200, updated_product.text
    approve_campaign_assets(client, auth, brand, product)

    payload = campaign_payload(brand, product)
    payload["create_default_items"] = True
    campaign_response = client.post(
        "/v1/campaign-packages",
        headers=auth,
        json=payload,
    )
    assert campaign_response.status_code == 201, campaign_response.text
    campaign = campaign_response.json()
    supply = create_approved_supply(client, auth, campaign, brand, product)

    for item in campaign["items"]:
        db.add(
            ContentVersion(
                organization_id=campaign["organization_id"],
                project_id=item["project"]["id"],
                brief_revision_id=campaign["current_brief_revision"]["id"],
                supply_snapshot_id=supply["id"],
                version_number=1,
                content={
                    "headline": "Verified seasonal produce",
                    "body": f"{product['name']} from Yunnan is ready for this campaign.",
                },
                status=ReviewStatus.approved,
                created_by=campaign["created_by"],
                reviewed_by=campaign["created_by"],
                review_note="Current evidence verified",
            )
        )
    db.commit()

    activated = client.patch(
        f"/v1/campaign-packages/{campaign['id']}/status",
        headers=auth,
        json={"status": "active"},
    )
    assert activated.status_code == 200, activated.text

    revised_product_payload = dict(initial_product_payload)
    revised_product_payload["origin"] = "Sichuan"
    revised_product = client.put(
        f"/v1/products/{product['id']}",
        headers=auth,
        json=revised_product_payload,
    )
    assert revised_product.status_code == 200, revised_product.text
    submitted_product = client.post(
        f"/v1/products/{product['id']}/submit",
        headers=auth,
    )
    assert submitted_product.status_code == 200, submitted_product.text
    reviewed_product = client.post(
        f"/v1/products/{product['id']}/review",
        headers=auth,
        json={"status": "approved", "note": "Revised product facts checked"},
    )
    assert reviewed_product.status_code == 200, reviewed_product.text

    stale_campaign = client.get(
        f"/v1/campaign-packages/{campaign['id']}",
        headers=auth,
    )
    assert stale_campaign.status_code == 200, stale_campaign.text
    assert all(item["approved_version_id"] is None for item in stale_campaign.json()["items"])
    assert all(
        "content_claims_unmapped" in item["stale_reasons"]
        for item in stale_campaign.json()["items"]
    )

    stale_download = client.get(
        f"/v1/campaign-packages/{campaign['id']}/presentation",
        headers=auth,
    )
    assert stale_download.status_code == 200, stale_download.text
    stale_deck = Presentation(BytesIO(stale_download.content))
    stale_text = "\n".join(
        shape.text
        for slide in stale_deck.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "请勿直接发布" in stale_text
